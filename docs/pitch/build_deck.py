#!/usr/bin/env python
"""Generate the 'Claros: the speech platform Windows already earned' pitch deck.

Reproducible build: `python build_deck.py` writes live-voiceover.pptx beside it.
Structure follows WHY -> WHAT -> PROOF -> HOW. Keep in sync with live-voiceover.md.

The value slide leaves a shape named "VIDEO_FRAME"; a review-only PowerPoint COM
step drops the actual mp4 onto it so the deck can play in a slideshow. The
committed pptx stays lightweight (poster placeholder only) so the ~120 MB demo
video is never committed.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --- palette (Microsoft-ish) ---
INK = RGBColor(0x1F, 0x1F, 0x1F)
BLUE = RGBColor(0x00, 0x78, 0xD4)
DEEP = RGBColor(0x10, 0x3A, 0x5E)
TEAL = RGBColor(0x00, 0x99, 0x8A)
AMBER = RGBColor(0xCA, 0x5F, 0x00)
GREY = RGBColor(0x60, 0x60, 0x60)
LIGHT = RGBColor(0xF3, 0xF6, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xEA, 0xF1, 0xFB)
SKY = RGBColor(0xCF, 0xE0, 0xF3)
STEEL = RGBColor(0x16, 0x4A, 0x74)

FONT = "Segoe UI"
FONT_L = "Segoe UI Light"
FONT_SB = "Segoe UI Semibold"

SW, SH = Inches(13.333), Inches(7.5)


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb


def para(tf, text, size, color=INK, bold=False, font=FONT, align=PP_ALIGN.LEFT,
         space_before=0, space_after=6, bullet=False, level=0, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    _bullet(p, bullet)
    return p


def _bullet(p, on):
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    if on:
        pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": FONT}))
        pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "\u2022"}))
    else:
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    _solid(sp, color)
    if line is not None:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def footer(slide, n):
    tb = box(slide, Inches(0.55), Inches(7.02), Inches(9), Inches(0.35))
    para(tb.text_frame, "Claros  \u00b7  one on-device speech platform, whole fleet",
         10, GREY, font=FONT, first=True, space_after=0)
    tb2 = box(slide, Inches(12.2), Inches(7.02), Inches(0.7), Inches(0.35))
    para(tb2.text_frame, str(n), 10, GREY, align=PP_ALIGN.RIGHT, first=True, space_after=0)


def accent_bar(slide, color=BLUE):
    rect(slide, 0, 0, SW, Inches(0.16), color)


def title_block(slide, kicker, title, color=BLUE):
    accent_bar(slide, color)
    tb = box(slide, Inches(0.55), Inches(0.45), Inches(12.2), Inches(0.4))
    para(tb.text_frame, kicker.upper(), 13, color, bold=True, font=FONT_SB,
         first=True, space_after=0)
    tt = box(slide, Inches(0.55), Inches(0.85), Inches(12.2), Inches(1.0))
    para(tt.text_frame, title, 30, INK, bold=True, font=FONT_SB, first=True, space_after=0)


def divider(prs, tag, title, sub=None, color=BLUE):
    s = blank(prs)
    rect(s, 0, 0, SW, SH, DEEP)
    rect(s, 0, 0, SW, Inches(0.22), color)
    rect(s, Inches(0.85), Inches(3.9), Inches(1.5), Inches(0.09), TEAL)
    tb = box(s, Inches(0.85), Inches(2.55), Inches(11.6), Inches(0.7))
    para(tb.text_frame, tag.upper(), 22, TEAL, bold=True, font=FONT_SB, first=True, space_after=0)
    tt = box(s, Inches(0.85), Inches(4.25), Inches(11.6), Inches(2.2))
    para(tt.text_frame, title, 38, WHITE, bold=True, font=FONT_SB, first=True, space_after=10)
    if sub:
        para(tt.text_frame, sub, 18, SKY, font=FONT_L, space_after=0)
    return s


def two_cards(slide, left, right, y=Inches(2.2), h=Inches(3.6)):
    """left/right = (heading, header_color, [lines...])."""
    for (head, col, lines), x in ((left, Inches(0.72)), (right, Inches(6.86))):
        rect(slide, x, y, Inches(5.75), h, WHITE, line=RGBColor(0xD9, 0xE4, 0xF0))
        rect(slide, x, y, Inches(5.75), Inches(0.62), col)
        tb = box(slide, x + Inches(0.32), y + Inches(0.12), Inches(5.1), Inches(0.5))
        para(tb.text_frame, head, 17, WHITE, bold=True, font=FONT_SB, first=True, space_after=0)
        bd = box(slide, x + Inches(0.32), y + Inches(0.85), Inches(5.1), h - Inches(1.0))
        for i, ln in enumerate(lines):
            para(bd.text_frame, ln, 14.5, INK, font=FONT, bullet=True,
                 first=(i == 0), space_after=10)


prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH

# ================================================================ 1 - TITLE + WHY hook (merged)
s = blank(prs)
rect(s, 0, 0, SW, SH, DEEP)
rect(s, 0, 0, SW, Inches(0.22), BLUE)
rect(s, 0, Inches(4.55), SW, Inches(0.05), TEAL)
tb = box(s, Inches(0.8), Inches(1.35), Inches(11.9), Inches(3.1))
para(tb.text_frame, "Claros", 46, WHITE, bold=True, font=FONT_SB,
     first=True, space_after=4)
para(tb.text_frame, "the speech platform Windows already earned", 38, TEAL, bold=True,
     font=FONT_SB, space_after=14)
para(tb.text_frame,
     "Apple quietly shipped a complete on-device speech platform to every Mac. "
     "Windows ships speech tech that is as good or better, and comparable to the "
     "cloud, then locks it away behind a missing, misaligned API that only runs on "
     "a sliver of the fleet. This is that platform: cohesive, whole-fleet, and "
     "proven today.",
     20, SKY, font=FONT_L, space_after=0)
tb2 = box(s, Inches(0.8), Inches(4.8), Inches(11.9), Inches(1.4))
para(tb2.text_frame,
     "One cohesive API for the whole Windows 11 fleet: enumerate voices, synthesize, "
     "recognize, converse, all on-device.",
     16, RGBColor(0x9F, 0xC4, 0xE7), font=FONT, first=True, space_after=2)
para(tb2.text_frame, "Local-first and free by default; the same code scales to Azure.",
     16, RGBColor(0x9F, 0xC4, 0xE7), bold=True, font=FONT_SB, space_after=2)
para(tb2.text_frame,
     "Claros is the working name of the implementation. It should ship as Windows.Speech.",
     16, TEAL, bold=True, font=FONT_SB, space_after=0)

# ================================================================ 2 - WHY: two Whispers, two APIs
s = blank(prs)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, "Why \u00b7 no cohesive platform",
            "Windows ships Whisper twice: two APIs, one engine", AMBER)
tb = box(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(0.7))
para(tb.text_frame,
     "Where Windows does push on-device speech, it ships the very same Whisper model "
     "behind two different, unrelated APIs, and neither is a platform:",
     16, INK, font=FONT, first=True, space_after=0)
two_cards(
    s,
    ("WINDOWS AI APIs", BLUE, [
        "On-device speech-to-text via the Windows AI / Windows App SDK surface.",
        "Its engine under the hood: Whisper (Large v3 Turbo).",
        "Its own object model, gated to Copilot+ PCs with an NPU.",
    ]),
    ("FOUNDRY LOCAL", DEEP, [
        "A separate local-model runtime and API, with its own SDK and CLI.",
        "Ships the same Whisper family again, a different way in.",
        "Same engine, different shape: two doors, no shared platform.",
    ]),
    y=Inches(2.75), h=Inches(2.75),
)
tb = box(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(0.9))
para(tb.text_frame,
     "Same model, two API shapes, and both inherit Whisper's real problem: it doesn't "
     "run well on the machines people actually have.",
     18, DEEP, bold=True, font=FONT_SB, first=True, space_after=0)
footer(s, 2)

# ================================================================ 3 - WHY: not optimized for Windows
s = blank(prs)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, "Why \u00b7 not optimized for the fleet",
            "Even ignore the API mess: the tech is too heavy", AMBER)
rect(s, Inches(0.7), Inches(2.0), Inches(5.7), Inches(3.35), RGBColor(0xFB, 0xF0, 0xE6))
tb = box(s, Inches(1.0), Inches(2.25), Inches(5.2), Inches(2.9))
para(tb.text_frame, "Whisper isn't built for Windows", 20, AMBER, bold=True, font=FONT_SB,
     first=True, space_after=8)
para(tb.text_frame,
     "The model that's being pushed simply isn't optimized for the platform. It's "
     "multiple gigabytes on disk and in RAM, and it runs slowly even on an NPU, so "
     "you need a very high-end PC, which we know most people do not run.",
     16, INK, font=FONT, space_after=0)
facts = [
    "Multiple GB on disk, just to install the speech model.",
    "Multiple GB of RAM at runtime, on a fleet where 16 GB is common.",
    "Slow even on an NPU; real-time is a struggle without top-tier silicon.",
    "So it only lands on a sliver of high-end, Copilot+ machines.",
    "Result: extremely limited adoption. Developers can't target it.",
]
tb = box(s, Inches(6.7), Inches(2.0), Inches(6.0), Inches(3.6))
for i, f in enumerate(facts):
    para(tb.text_frame, f, 16.5, INK, font=FONT, bullet=True, first=(i == 0), space_after=13)
tb = box(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.9))
para(tb.text_frame,
     "A platform only counts if it runs on the fleet you actually have. Today's path "
     "doesn't, on hardware, on disk, or on memory.",
     18, DEEP, bold=True, font=FONT_SB, first=True, space_after=0)
footer(s, 3)

# ================================================================ 4 - WHY: Apple shipped it
s = blank(prs)
rect(s, 0, 0, SW, SH, LIGHT)
title_block(s, "Why \u00b7 Apple already did it",
            "Apple quietly shipped the complete platform, to every Mac", TEAL)
tb = box(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(0.7))
para(tb.text_frame,
     "While Windows debated, Apple shipped a cohesive on-device speech platform and "
     "put it on every Mac in use, not just the newest:",
     16, INK, font=FONT, first=True, space_after=0)
tiles = [
    ("SpeechAnalyzer (STT)", "On-device recognition with ITN, punctuation, per-speaker attribution.", TEAL),
    ("Neural TTS", "High-quality on-device voices, one shared system model.", TEAL),
    ("Runs on every Mac", "Not gated to the latest silicon; the whole fleet is the target.", BLUE),
    ("~900 MB, real-time", "Efficient enough to be practical: ships in real apps today.", BLUE),
]
positions = [(0.72, 2.75), (6.86, 2.75), (0.72, 4.55), (6.86, 4.55)]
for (tx, ty), (h, d, col) in zip(positions, tiles):
    rect(s, Inches(tx), Inches(ty), Inches(5.75), Inches(1.55), WHITE, line=RGBColor(0xDD, 0xE6, 0xF1))
    rect(s, Inches(tx), Inches(ty), Inches(0.12), Inches(1.55), col)
    tb = box(s, Inches(tx + 0.35), Inches(ty + 0.22), Inches(5.15), Inches(1.15))
    para(tb.text_frame, h, 17, INK, bold=True, font=FONT_SB, first=True, space_after=6)
    para(tb.text_frame, d, 14, GREY, font=FONT, space_after=0)
tb = box(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.6))
para(tb.text_frame,
     "The model is proven and wanted: cohesive, on-device, fleet-wide speech. The "
     "only question is whether Windows will answer it.",
     16, DEEP, bold=True, font=FONT_SB, first=True, space_after=0)
footer(s, 4)

# ================================================================ 5 - WHY: the face-palm
s = blank(prs)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, "Why \u00b7 the face-palm",
            "Windows already ships tech that matches or beats it", AMBER)
tb = box(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(0.9))
para(tb.text_frame,
     "Here is the frustrating part. Windows doesn't need to invent anything. The "
     "tech already ships in the box, and it's as good or better than the Mac, with "
     "results comparable to the cloud:",
     16, INK, font=FONT, first=True, space_after=0)
rows = [
    ("Natural HD voices", "The same neural voices as Microsoft's cloud, running offline on the device."),
    ("Live Captions recognizer", "On-device STT that runs great on the CPU, at ~500 MB, no NPU required."),
    ("Efficient by design", "Light on disk and RAM, fast on an ordinary CPU, and better with an NPU."),
]
y = 2.95
for i, (a, b) in enumerate(rows):
    rect(s, Inches(0.7), Inches(y + i * 1.0), Inches(0.12), Inches(0.82), BLUE)
    tbb = box(s, Inches(1.0), Inches(y + i * 1.0), Inches(11.4), Inches(0.82))
    para(tbb.text_frame, a, 19, INK, bold=True, font=FONT_SB, first=True, space_after=2)
    para(tbb.text_frame, b, 15, GREY, font=FONT, space_after=0)
rect(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.9), RGBColor(0xFB, 0xF0, 0xE6))
tb = box(s, Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.65))
tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tb.text_frame,
     "Microsoft has the winning tech on the device, and simply doesn't hand it to "
     "developers. That is the whole problem, and it is entirely fixable.",
     16, AMBER, bold=True, font=FONT_SB, first=True, space_after=0)
footer(s, 5)

# ================================================================ 6 - WHAT divider
divider(prs, "What",
        "Claros: one platform, the whole fleet",
        "A complete, cohesive speech API that runs beautifully on every Windows 11 "
        "PC, Copilot+ or not, because the models are that efficient.", color=TEAL)

# ================================================================ 7 - WHAT: introduce Claros
s = blank(prs)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, "What \u00b7 the platform",
            "One cohesive API: speak, listen, converse", BLUE)
tb = box(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.7))
para(tb.text_frame,
     "Claros unifies the scattered pieces into a single, consistent surface, "
     "and it targets the fleet, not just Copilot+ PCs:",
     16, INK, font=FONT, first=True, space_after=0)
tiles = [
    ("Speak", "Enumerate installed Natural HD voices and synthesize offline, live, with word-boundary events.", TEAL),
    ("Listen", "Recognize on-device with the Live Captions engine, one recognizer per source, clean finals.", BLUE),
    ("Converse", "A full round-trip loop, capture to reply, with barge-in and a place to plug in intelligence.", DEEP),
    ("Fleet-wide & efficient", "Runs great on CPU, better on NPU. Local-first and free; same code scales to Azure.", AMBER),
]
positions = [(0.72, 2.65), (6.86, 2.65), (0.72, 4.5), (6.86, 4.5)]
for (tx, ty), (h, d, col) in zip(positions, tiles):
    rect(s, Inches(tx), Inches(ty), Inches(5.75), Inches(1.65), WHITE, line=RGBColor(0xDD, 0xE6, 0xF1))
    rect(s, Inches(tx), Inches(ty), Inches(0.12), Inches(1.65), col)
    tb = box(s, Inches(tx + 0.35), Inches(ty + 0.24), Inches(5.15), Inches(1.2))
    para(tb.text_frame, h, 19, INK, bold=True, font=FONT_SB, first=True, space_after=6)
    para(tb.text_frame, d, 14, GREY, font=FONT, space_after=0)
tb = box(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.6))
para(tb.text_frame,
     "One name, one shape, whole fleet. Claros is what \u2018Windows.Speech\u2019 should "
     "be: a platform a developer can just pick up, and their customers can actually run.",
     16, DEEP, bold=True, font=FONT_SB, first=True, space_after=0)
footer(s, 7)

# ================================================================ 8 - WHAT: the value it unlocks (demo)
s = blank(prs)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, "What it enables \u00b7 the payoff",
            "Amazing quality, generated in real time", TEAL)
# video frame (COM drops the mp4 here for review; poster placeholder otherwise)
vf = rect(s, Inches(0.7), Inches(1.95), Inches(7.7), Inches(4.35), RGBColor(0x0B, 0x24, 0x3B))
vf.name = "VIDEO_FRAME"
pb = rect(s, Inches(4.15), Inches(3.75), Inches(0.8), Inches(0.8), TEAL, MSO_SHAPE.OVAL)
pb.name = "VIDEO_PLAY"
tri = rect(s, Inches(4.42), Inches(3.95), Inches(0.32), Inches(0.4), WHITE, MSO_SHAPE.ISOSCELES_TRIANGLE)
tri.rotation = 90
tri.name = "VIDEO_PLAY_TRI"
cap = box(s, Inches(0.9), Inches(5.75), Inches(7.3), Inches(0.4))
para(cap.text_frame, "WinUI sample \u00b7 live on-device voiceover from a subtitle track",
     13, SKY, font=FONT, first=True, space_after=0)
# side narrative
tb = box(s, Inches(8.75), Inches(2.0), Inches(3.95), Inches(4.3))
para(tb.text_frame,
     "This is the value it unlocks: studio-grade speech, generated live on the "
     "device, for free.",
     16, DEEP, bold=True, font=FONT_SB, first=True, space_after=12)
for t in [
    "Indistinguishable from cloud neural voices, running fully offline.",
    "Generated in real time on an ordinary CPU, no NPU, no wait.",
    "Delivers the modern Microsoft mission: unmetered intelligence on every desk and in every home.",
]:
    para(tb.text_frame, t, 14.5, INK, font=FONT, bullet=True, space_after=12)
footer(s, 8)

# ================================================================ 9 - WHAT: PowerPoint voiceover, zero cloud
s = blank(prs)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, "What it enables \u00b7 PowerPoint",
            "Narrate every deck, with zero cloud cost", TEAL)
two_cards(
    s,
    ("POWERPOINT + CLAROS", TEAL, [
        "Speaker notes become the narration script, one per slide.",
        "Render the deck to a fully narrated MP4, on the device.",
        "One deck, many languages, no re-recording.",
        "Zero cloud cost: nothing is metered, nothing leaves the machine.",
    ]),
    ("CLOUD VIDEO TOOLS (e.g. Synthesia)", GREY, [
        "Per-minute or subscription cloud cost.",
        "Script and slides leave the org to render.",
        "Online-only; a separate tool and export step.",
        "No on-device, private, or free tier.",
    ]),
    y=Inches(2.15), h=Inches(3.45),
)
tb = box(s, Inches(0.72), Inches(5.95), Inches(11.9), Inches(0.8))
para(tb.text_frame,
     "Clipchamp and PowerPoint voiceover run on Azure today; Claros makes the "
     "same narration local and free. Even this deck could narrate itself, on-device.",
     16, DEEP, bold=True, font=FONT_SB, first=True, space_after=0)
footer(s, 9)

# ================================================================ 10 - PROOF divider
divider(prs, "Proof",
        "Efficient enough to change what teams can build",
        "The clearest proof isn't a benchmark, it's a real product that couldn't ship "
        "until the speech got this light.", color=BLUE)

# ================================================================ 11 - PROOF: Contoso Finance
s = blank(prs)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, "Proof \u00b7 Contoso Finance",
            "From \u2018barely runs\u2019 to \u2018room for a local LLM\u2019", BLUE)
tb = box(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.65))
para(tb.text_frame,
     "Contoso Finance transcribes two-party advisor calls on the analyst's own PC. "
     "On a 16 GB machine, the speech engine decided whether the app could exist:",
     16, INK, font=FONT, first=True, space_after=0)
rect(s, Inches(0.72), Inches(2.7), Inches(5.75), Inches(2.3), RGBColor(0xFB, 0xF0, 0xE6))
rect(s, Inches(0.72), Inches(2.7), Inches(5.75), Inches(0.6), AMBER)
tb = box(s, Inches(1.02), Inches(2.81), Inches(5.15), Inches(0.5))
para(tb.text_frame, "BEFORE \u00b7 Whisper (NPU or CPU)", 15, WHITE, bold=True, font=FONT_SB, first=True, space_after=0)
tb = box(s, Inches(1.02), Inches(3.45), Inches(5.15), Inches(1.45))
para(tb.text_frame, "\u2248 4 GB", 38, AMBER, bold=True, font=FONT_SB, first=True, space_after=2)
para(tb.text_frame,
     "just for speech, on a 16 GB PC. Almost no headroom left, and ~4\u00d7 the memory "
     "of the Mac doing the same job.",
     14.5, INK, font=FONT, space_after=0)
rect(s, Inches(6.86), Inches(2.7), Inches(5.75), Inches(2.3), RGBColor(0xE4, 0xF3, 0xF0))
rect(s, Inches(6.86), Inches(2.7), Inches(5.75), Inches(0.6), TEAL)
tb = box(s, Inches(7.16), Inches(2.81), Inches(5.15), Inches(0.5))
para(tb.text_frame, "AFTER \u00b7 Claros", 15, WHITE, bold=True, font=FONT_SB, first=True, space_after=0)
tb = box(s, Inches(7.16), Inches(3.45), Inches(5.15), Inches(1.45))
para(tb.text_frame, "\u2248 500 MB", 38, TEAL, bold=True, font=FONT_SB, first=True, space_after=2)
para(tb.text_frame,
     "same transcription, same accuracy, at roughly the Mac's footprint, freeing "
     "gigabytes back to the machine.",
     14.5, INK, font=FONT, space_after=0)
rect(s, Inches(6.4), Inches(3.7), Inches(0.55), Inches(0.4), BLUE, MSO_SHAPE.RIGHT_ARROW)
rect(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.25), CARD)
tb = box(s, Inches(1.0), Inches(5.48), Inches(11.3), Inches(0.9))
tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tb.text_frame,
     "That freed memory does two things: it lets the app ship on a normal 16 GB PC at "
     "all, and it leaves room to run a local LLM alongside the transcription, "
     "impossible today under Whisper's memory pressure. Efficiency is what makes "
     "on-device intelligence practical.",
     16, DEEP, bold=True, font=FONT_SB, first=True, space_after=0)
footer(s, 11)

# ================================================================ 12 - PROOF: STT benchmark table
s = blank(prs)
rect(s, 0, 0, SW, SH, WHITE)
title_block(s, "Proof \u00b7 the numbers",
            "Mac-class accuracy at a fraction of the memory", BLUE)
tb = box(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(0.5))
para(tb.text_frame,
     "Real 58 s two-party mortgage call, normalized to a 2-leg call (one recognizer "
     "per speaker); single-stream engines doubled.",
     13.5, GREY, font=FONT, first=True, space_after=0)
cols = [(0.7, 3.7), (4.5, 1.35), (5.95, 2.0), (8.05, 1.55), (9.65, 2.95)]
heads = ["Engine", "First final", "Peak RAM (2 legs)", "Hardware", "Numbers / ITN"]
hy = 2.2
rect(s, Inches(0.7), Inches(hy), Inches(11.9), Inches(0.5), DEEP)
for (cx, cw), h in zip(cols, heads):
    hb = box(s, Inches(cx + 0.1), Inches(hy), Inches(cw - 0.15), Inches(0.5))
    hb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(hb.text_frame, h, 12.5, WHITE, bold=True, font=FONT_SB, first=True, space_after=0)
rows = [
    ("Claros (Live Captions)", "~4 s", "~500 MB", "CPU", "Words yes; ITN off today (fixable)", "ours"),
    ("Apple SpeechAnalyzer (macOS)", "~4 s", "~440 MB", "Apple ANE", "Full ITN: $610,000, 6.2%", "peer"),
    ("WinAI Speech Preview", "3.5 s", "~6,400 MB", "Hexagon NPU", "Full ITN (Whisper Turbo)", "peer"),
    ("Nemotron 0.6B (Foundry Local)", "1.2 s", "~1,750 MB", "CPU", "No clean sentence breaks", "out"),
    ("Parakeet TDT 0.6b", "1.4 s", "~1,780 MB / leg", "CPU", "Good ITN, but duplicates finals", "out"),
    ("Whisper small (CPU ONNX)", "2.3 s", "~1,200 MB", "CPU", "Low: hallucinates numbers", "out"),
]
ry = hy + 0.5
for cells in rows:
    tag = cells[5]
    bg = RGBColor(0xE4, 0xF3, 0xF0) if tag == "ours" else (LIGHT if tag == "peer" else WHITE)
    rect(s, Inches(0.7), Inches(ry), Inches(11.9), Inches(0.56), bg)
    if tag == "ours":
        rect(s, Inches(0.7), Inches(ry), Inches(0.1), Inches(0.56), TEAL)
    for (cx, cw), val, idx in zip(cols, cells[:5], range(5)):
        cb = box(s, Inches(cx + 0.1), Inches(ry), Inches(cw - 0.15), Inches(0.56))
        cb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        col = INK if idx == 0 else GREY
        bold = (idx == 0)
        if idx == 2 and tag in ("ours", "peer"):
            col, bold = BLUE, True
        para(cb.text_frame, val, 12.5, col, bold=bold,
             font=FONT_SB if bold else FONT, first=True, space_after=0)
    ry += 0.56
rect(s, Inches(0.7), Inches(ry + 0.12), Inches(11.9), Inches(1.05), CARD)
tb = box(s, Inches(1.0), Inches(ry + 0.24), Inches(11.3), Inches(0.85))
tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tb.text_frame,
     "Claros matches the peers on word accuracy at ~500 MB for two legs, "
     "~13\u00d7 less RAM than the NPU path, no NPU. One honest gap: it emits spelled-out "
     "numbers today because the shipping ITN model sits behind a native finalizer we "
     "disable for ARM64 stability, fixable in-box. The rest are out on quality or, for "
     "Parakeet, on being unable to hold real time once finals must be immutable.",
     13.5, DEEP, bold=True, font=FONT_SB, first=True, space_after=0)
footer(s, 12)

# ================================================================ 13 - HOW divider
divider(prs, "How",
        "Make this implementation real",
        "It works today. The only thing between this hack and a real platform is a "
        "decision Microsoft already owns.", color=AMBER)

# ================================================================ 14 - HOW: local-first, Azure-optional
s = blank(prs)
rect(s, 0, 0, SW, SH, LIGHT)
title_block(s, "How \u00b7 the design", "Local by default. The same SDK scales to Azure.", BLUE)
rect(s, Inches(0.72), Inches(2.05), Inches(5.75), Inches(3.0), WHITE, line=RGBColor(0xD9, 0xE4, 0xF0))
rect(s, Inches(0.72), Inches(2.05), Inches(5.75), Inches(0.62), TEAL)
tb = box(s, Inches(1.04), Inches(2.17), Inches(5.1), Inches(0.5))
para(tb.text_frame, "DEFAULT: on-device", 16, WHITE, bold=True, font=FONT_SB, first=True, space_after=0)
tb = box(s, Inches(1.04), Inches(2.9), Inches(5.1), Inches(2.0))
for i, t in enumerate([
    "Offline, private, free, instant.",
    "Runs on the whole fleet, no NPU required.",
    "EmbeddedSpeechConfig on the installed voice + Live Captions model.",
]):
    para(tb.text_frame, t, 15, INK, font=FONT, bullet=True, first=(i == 0), space_after=9)
rect(s, Inches(6.86), Inches(2.05), Inches(5.75), Inches(3.0), WHITE, line=RGBColor(0xD9, 0xE4, 0xF0))
rect(s, Inches(6.86), Inches(2.05), Inches(5.75), Inches(0.62), BLUE)
tb = box(s, Inches(7.18), Inches(2.17), Inches(5.1), Inches(0.5))
para(tb.text_frame, "UPSELL: same code + your Azure creds", 16, WHITE, bold=True, font=FONT_SB, first=True, space_after=0)
tb = box(s, Inches(7.18), Inches(2.9), Inches(5.1), Inches(2.0))
for i, t in enumerate([
    "Swap in a cloud SpeechConfig(key, region).",
    "Same SpeechSynthesizer / SpeechRecognizer calls, same voices.",
    "More voices, more languages, server-side scale.",
]):
    para(tb.text_frame, t, 15, INK, font=FONT, bullet=True, first=(i == 0), space_after=9)
rect(s, Inches(6.4), Inches(3.35), Inches(0.55), Inches(0.4), AMBER, MSO_SHAPE.RIGHT_ARROW)
rect(s, Inches(0.72), Inches(5.35), Inches(11.9), Inches(1.05), CARD)
tb = box(s, Inches(1.02), Inches(5.5), Inches(11.3), Inches(0.8))
tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tb.text_frame,
     "This isn't hand-waving: the reference implementation is the Azure Speech SDK in "
     "embedded mode. The upgrade to Azure is a config swap, not a rewrite. Local is "
     "the free on-ramp; Azure plus Microsoft's own apps is a funnel no rival can match.",
     15, DEEP, bold=True, font=FONT_SB, first=True, space_after=0)
footer(s, 14)

# ================================================================ 15 - HOW: the ask / CTA
s = blank(prs)
rect(s, 0, 0, SW, SH, DEEP)
rect(s, 0, 0, SW, Inches(0.16), BLUE)
tb = box(s, Inches(0.7), Inches(0.6), Inches(12), Inches(1.2))
para(tb.text_frame, "THE ASK TO MICROSOFT", 14, RGBColor(0x8F, 0xC4, 0xF0), bold=True,
     font=FONT_SB, first=True, space_after=6)
para(tb.text_frame, "Make this implementation real", 32, WHITE, bold=True,
     font=FONT_SB, space_after=0)
tb = box(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(1.0))
para(tb.text_frame,
     "This repo is a complete, working reference implementation. It runs today only "
     "by hacking around the on-device licensing, and that hack was trivial to bypass. "
     "There is no technical barrier left, only a decision:",
     16, SKY, font=FONT, first=True, space_after=0)
rect(s, Inches(0.7), Inches(3.05), Inches(11.9), Inches(1.05), STEEL)
tb = box(s, Inches(1.0), Inches(3.2), Inches(11.3), Inches(0.75))
tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tb.text_frame,
     "The only blocker is a licensing key, and Microsoft controls it.",
     22, WHITE, bold=True, font=FONT_SB, first=True, space_after=0)
bl = [
    "Bless the license so any app can reach the on-device runtime, no hack.",
    "Unite the scattered efforts (WinAI, Foundry, HD voices, Live Captions) behind one platform.",
    "Ship it as Windows.Speech: cohesive, fleet-wide, on-device, with a paved road to Azure.",
]
tb = box(s, Inches(1.0), Inches(4.4), Inches(11.4), Inches(1.6))
for i, t in enumerate(bl):
    para(tb.text_frame, t, 16, WHITE, font=FONT, bullet=True, first=(i == 0), space_after=9)
tb = box(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.7))
para(tb.text_frame,
     "The tech is done, and Claros proves it. Get the company together, unite behind "
     "one platform, and compete with macOS: ship it as Windows.Speech.",
     16, TEAL, bold=True, font=FONT_SB, first=True, space_after=0)

out = Path(__file__).with_name("live-voiceover.pptx")
prs.save(out)
print(f"wrote {out}  ({out.stat().st_size:,} bytes, {len(prs.slides._sldIdLst)} slides)")
